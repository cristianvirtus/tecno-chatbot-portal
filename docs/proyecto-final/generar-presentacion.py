"""Genera nodi-presentacion.pptx (10 diapositivas, 16:9).

Uso: python3 docs/proyecto-final/generar-presentacion.py
Requiere: python-pptx y rsvg-convert (los íconos se rasterizan desde SVG).
"""

import math
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Paleta de marca del curso (extraída del CSS de datapath.ai).
NAVY = "05192D"
NAVY_DEEP = "020F1C"
CARD = "0A2540"
CARD_ALT = "0D2B45"
LINE = "1A3A55"
LINE_SOFT = "2A4A65"
ORANGE = "FB5604"
GREEN = "00C48C"
PURPLE = "3D348B"
WHITE = "FFFFFF"
MUTED = "7A9BB5"

FONT = "Calibri"
W, H = 13.333, 7.5
MARGIN = 0.62
CONTENT_W = W - 2 * MARGIN

ICONS = {
    "user": '<circle cx="12" cy="8" r="4"/><path d="M4 21a8 8 0 0 1 16 0"/>',
    "browser": '<rect x="3" y="4" width="18" height="16" rx="2"/><path d="M3 9h18M6.5 6.5h.01M9.5 6.5h.01"/>',
    "code": '<path d="m8 7-5 5 5 5m8-10 5 5-5 5m-3-13-2 16"/>',
    "model": (
        '<path d="M9 4a3 3 0 0 0-3 3v1a4 4 0 0 0-1 7 3 3 0 0 0 4 4M15 4a3 3 0 0 1 3 3v1'
        'a4 4 0 0 1 1 7 3 3 0 0 1-4 4M12 3v18"/><path d="M9 9h3M12 15h3"/>'
    ),
    "tool": '<path d="M14 6a4 4 0 0 0-5 5L4 16l4 4 5-5a4 4 0 0 0 5-5l-3 3-4-4Z"/>',
    "database": (
        '<ellipse cx="12" cy="5" rx="8" ry="3"/>'
        '<path d="M4 5v6c0 1.7 3.6 3 8 3s8-1.3 8-3V5M4 11v6c0 1.7 3.6 3 8 3s8-1.3 8-3v-6"/>'
    ),
    "search": '<circle cx="10.5" cy="10.5" r="6.5"/><path d="m15.5 15.5 5 5"/>',
    "stream": '<path d="M4 7h10M4 12h16M4 17h12"/><path d="m17 5 3 2-3 2"/>',
    "spark": (
        '<path d="m12 3 1.6 5.4L19 10l-5.4 1.6L12 17l-1.6-5.4L5 10l5.4-1.6Z"/>'
        '<path d="m18.5 16 .7 2.3 2.3.7-2.3.7-.7 2.3-.7-2.3-2.3-.7 2.3-.7Z"/>'
    ),
    "shield": '<path d="M12 3 5 6v5c0 5 3 8 7 10 4-2 7-5 7-10V6Z"/><path d="m9 12 2 2 4-4"/>',
    "mic": '<rect x="9" y="3" width="6" height="12" rx="3"/><path d="M5 11a7 7 0 0 0 14 0M12 18v3M9 21h6"/>',
    "git": (
        '<circle cx="6" cy="5" r="2"/><circle cx="6" cy="19" r="2"/><circle cx="18" cy="12" r="2"/>'
        '<path d="M6 7v10M8 6c0 3.5 2 6 8 6"/>'
    ),
    "pipeline": (
        '<rect x="3" y="4" width="5" height="5" rx="1"/><rect x="16" y="4" width="5" height="5" rx="1"/>'
        '<rect x="9.5" y="15" width="5" height="5" rx="1"/>'
        '<path d="M8 6.5h8M5.5 9v3l6.5 3m6.5-6v3L12 15"/>'
    ),
    "deploy": '<path d="M12 3v12m-4-8 4-4 4 4"/><path d="M5 14v5h14v-5"/>',
    "check": '<path d="m4 12 5 5L20 6"/>',
    "key": '<circle cx="7" cy="16" r="3"/><path d="m9.5 13.5 9-9m-3 3 3 3m-6 0 2 2"/>',
    "agent": (
        '<rect x="4" y="6" width="16" height="13" rx="4"/>'
        '<path d="M9 11h.01M15 11h.01M9.5 15c1.7 1.1 3.3 1.1 5 0M12 6V3M9.5 3h5"/>'
    ),
    "play": '<circle cx="12" cy="12" r="9"/><path d="m10 8 6 4-6 4Z"/>',
    "cloud": '<path d="M6.5 18h11a4 4 0 0 0 .6-8 6.3 6.3 0 0 0-12-1.4A4.8 4.8 0 0 0 6.5 18Z"/>',
    "layers": '<path d="m12 3 9 5-9 5-9-5Z"/><path d="m3 13 9 5 9-5M3 17l9 5 9-5"/>',
}

SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24" width="{px}" height="{px}">'
    '<g fill="none" stroke="{color}" stroke-width="{sw}" stroke-linecap="round" '
    'stroke-linejoin="round">{inner}</g></svg>'
)

_tmp = Path(tempfile.mkdtemp(prefix="nodi-icons-"))
_cache: dict[tuple, Path] = {}


def icon(name: str, color: str, px: int = 128, sw: float = 1.8) -> Path:
    key = (name, color, px, sw)
    if key not in _cache:
        svg = _tmp / f"{name}-{color}-{px}-{sw}.svg"
        png = svg.with_suffix(".png")
        svg.write_text(SVG.format(px=px, color=f"#{color}", inner=ICONS[name], sw=sw))
        subprocess.run(
            ["rsvg-convert", "-w", str(px), "-h", str(px), str(svg), "-o", str(png)],
            check=True,
        )
        _cache[key] = png
    return _cache[key]


def rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def rect(slide, x, y, w, h, fill=None, line=None, radius=0.07, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        # El adjustment es relativo al lado corto de la figura.
        s.adjustments[0] = min(0.5, radius / min(w, h)) if min(w, h) else 0
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    s.text_frame.text = ""
    return s


def text(
    slide,
    x,
    y,
    w,
    h,
    body,
    size=14,
    color=WHITE,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    spacing=1.0,
    caps=False,
    space=0.0,
    italic=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = body if isinstance(body, list) else [body]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line.upper() if caps else line
        f = run.font
        f.name = FONT
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = rgb(color)
        if space:
            # El espaciado entre letras no está en la API pública de python-pptx.
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(space * 100)))
    return box


def place_icon(slide, name, x, y, size, color):
    slide.shapes.add_picture(str(icon(name, color)), Inches(x), Inches(y), Inches(size), Inches(size))


def icon_tile(slide, name, x, y, size=0.52, bg=ORANGE, fg=WHITE, radius=0.1):
    rect(slide, x, y, size, size, fill=bg, radius=radius)
    pad = size * 0.24
    place_icon(slide, name, x + pad, y + pad, size - 2 * pad, fg)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    shape.shadow.inherit = False


def flow_arrow(slide, x, y_center, length, color=ORANGE, left=False):
    """Línea horizontal con punta triangular."""
    head = min(0.13, length * 0.5)
    bar_x = x + head if left else x
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(bar_x), Inches(y_center - 0.011), Inches(length - head), Inches(0.022)
    )
    _solid(bar, color)
    tip = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(x if left else x + length - head),
        Inches(y_center - head / 2),
        Inches(head),
        Inches(head),
    )
    tip.rotation = 270 if left else 90
    _solid(tip, color)


def down_arrow(slide, x_center, y, length, color=ORANGE):
    head = min(0.13, length * 0.5)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x_center - 0.011), Inches(y), Inches(0.022), Inches(length - head)
    )
    _solid(bar, color)
    tip = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(x_center - head / 2),
        Inches(y + length - head),
        Inches(head),
        Inches(head),
    )
    tip.rotation = 180
    _solid(tip, color)


def pill(slide, x, y, label, fill=None, border=LINE_SOFT, color=WHITE, size=11, height=0.3):
    w = 0.2 + len(label) * size * 0.0092
    s = rect(slide, x, y, w, height, fill=fill, line=None if fill else border, radius=height / 2)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = rgb(color)
    return w


def new_slide(prs, bg=NAVY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(bg)
    return slide


def header(slide, kicker, title, subtitle=None, number=None, total=10):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(0.085))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(ORANGE)
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(slide, MARGIN, 0.4, 8, 0.22, kicker, size=10.5, color=ORANGE, bold=True, caps=True, space=1.6)
    text(slide, MARGIN, 0.66, CONTENT_W - 1.0, 0.55, title, size=31, bold=True, spacing=0.95)
    y = 1.28
    if subtitle:
        text(slide, MARGIN, y, CONTENT_W - 1.4, 0.32, subtitle, size=13.5, color=MUTED)
        y += 0.42
    if number:
        text(
            slide,
            W - MARGIN - 1.2,
            0.42,
            1.2,
            0.24,
            f"{number:02d} / {total:02d}",
            size=11,
            color=LINE_SOFT,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
    return y


def footer(slide, label="Nodi · Proyecto final"):
    text(slide, MARGIN, H - 0.44, 6, 0.22, label, size=9.5, color=LINE_SOFT)


# --------------------------------------------------------------------------- #
# Diapositivas
# --------------------------------------------------------------------------- #


def slide_cover(prs):
    slide = new_slide(prs, NAVY_DEEP)
    rect(slide, 0, 0, 0.16, H, fill=ORANGE, radius=0)
    text(slide, 1.0, 1.05, 7.4, 0.24, "Proyecto final · AI Agentic Engineer", size=12, color=ORANGE, bold=True, caps=True, space=1.4)
    text(slide, 0.96, 1.55, 7.4, 1.5, "Nodi", size=96, bold=True, spacing=0.85)
    text(slide, 1.0, 3.05, 7.2, 0.5, "Mentor de tecnología con IA", size=30, color=GREEN, bold=True)
    text(
        slide,
        1.0,
        3.75,
        6.7,
        0.9,
        "Un mini-agente educativo que conversa por texto o voz, consulta conocimiento real "
        "y muestra visualmente cómo construye cada respuesta.",
        size=14.5,
        color=MUTED,
        spacing=1.25,
    )
    x = 1.0
    for label in ["1 herramienta", "RAG local", "Streaming", "CI/CD"]:
        x += pill(slide, x, 4.95, label, fill=None, border=LINE_SOFT, color=WHITE) + 0.14
    text(slide, 1.0, 6.05, 6, 0.24, "Cristian Inda", size=13, bold=True)
    text(slide, 1.0, 6.32, 6, 0.24, "tecno-chatbot-portal.vercel.app", size=11.5, color=MUTED)

    # Constelación: agente al centro, seis capacidades alrededor.
    cx, cy, r = 10.55, 3.75, 1.95
    nodes = [("user", "Pregunta"), ("mic", "Voz"), ("model", "Claude"), ("database", "Corpus"), ("stream", "Stream"), ("pipeline", "CI/CD")]
    pts = []
    for i in range(6):
        a = math.radians(-90 + i * 60)
        pts.append((cx + r * math.cos(a), cy + r * math.sin(a)))
    for px, py in pts:
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cx), Inches(cy), Inches(px), Inches(py))
        c.line.color.rgb = rgb(LINE)
        c.line.width = Pt(1.25)
    for (px, py), (name, label) in zip(pts, nodes):
        d = 0.86
        rect(slide, px - d / 2, py - d / 2, d, d, fill=CARD, line=LINE_SOFT, radius=d / 2)
        place_icon(slide, name, px - 0.21, py - 0.21, 0.42, GREEN)
        text(slide, px - 0.9, py + 0.5, 1.8, 0.22, label, size=9.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    d = 1.65
    rect(slide, cx - d / 2, cy - d / 2, d, d, fill=ORANGE, radius=d / 2)
    place_icon(slide, "agent", cx - 0.42, cy - 0.42, 0.84, WHITE)
    return slide


def slide_agenda(prs):
    slide = new_slide(prs)
    header(slide, "Ruta de la presentación", "Agenda", "Del problema al producto desplegado", number=2)
    items = [
        ("01", "Reto y solución", "Por qué existe Nodi", "spark"),
        ("02", "Demo end-to-end", "La experiencia en vivo", "play"),
        ("03", "Arquitectura", "Qué ocurre detrás de cada respuesta", "layers"),
        ("04", "Inteligencia", "Tool use, prompt y RAG", "model"),
        ("05", "Calidad", "Pruebas automatizadas", "shield"),
        ("06", "DevOps", "De GitHub a Vercel", "pipeline"),
    ]
    cw, ch, gap = (CONTENT_W - 0.3) / 2, 0.82, 0.22
    for i, (n, title, detail, ic) in enumerate(items):
        col, row = i % 2, i // 2
        x = MARGIN + col * (cw + 0.3)
        y = 2.05 + row * (ch + gap)
        rect(slide, x, y, cw, ch, fill=CARD, line=LINE, radius=0.1)
        text(slide, x + 0.28, y + 0.17, 0.6, 0.4, n, size=22, color=ORANGE, bold=True)
        text(slide, x + 0.95, y + 0.16, cw - 2.0, 0.28, title, size=15, bold=True)
        text(slide, x + 0.95, y + 0.46, cw - 2.0, 0.26, detail, size=11.5, color=MUTED)
        icon_tile(slide, ic, x + cw - 0.78, y + 0.17, size=0.48, bg=LINE, fg=GREEN)

    y = 2.05 + 3 * (ch + gap) + 0.12
    rect(slide, MARGIN, y, CONTENT_W, 0.92, fill=CARD_ALT, line=ORANGE, radius=0.1)
    text(slide, MARGIN + 0.3, y + 0.17, 3.2, 0.24, "Idea conductora", size=11, color=ORANGE, bold=True, caps=True, space=1.2)
    text(slide, MARGIN + 0.3, y + 0.44, CONTENT_W - 0.6, 0.34, "Una función · un corpus · un stream · un push", size=19, bold=True)
    footer(slide)
    return slide


def slide_challenge(prs):
    slide = new_slide(prs)
    header(slide, "01 · Reto y solución", "Aprender tecnología también puede ser visual", "El reto: convertir una respuesta de IA en una experiencia didáctica", number=3)

    cw = (CONTENT_W - 0.95) / 2
    y = 2.15
    rect(slide, MARGIN, y, cw, 3.85, fill=CARD, line=LINE, radius=0.12)
    icon_tile(slide, "user", MARGIN + 0.35, y + 0.35, size=0.6, bg=LINE, fg=MUTED)
    text(slide, MARGIN + 0.35, y + 1.12, cw - 0.7, 0.3, "El problema", size=19, bold=True)
    text(
        slide,
        MARGIN + 0.35,
        y + 1.55,
        cw - 0.7,
        1.0,
        "Los conceptos técnicos llegan como bloques de texto: sin recorrido visual "
        "y sin saber de dónde salió la información.",
        size=13,
        color=MUTED,
        spacing=1.3,
    )
    rect(slide, MARGIN + 0.35, y + 2.7, cw - 0.7, 0.85, fill=NAVY_DEEP, line=None, radius=0.1)
    text(slide, MARGIN + 0.55, y + 2.86, cw - 1.1, 0.22, "Fricción", size=9.5, color=ORANGE, bold=True, caps=True, space=1.2)
    text(slide, MARGIN + 0.55, y + 3.11, cw - 1.1, 0.3, "“Entiendo la definición, pero no veo el flujo.”", size=12.5, italic=True)

    ax = MARGIN + cw + 0.13
    flow_arrow(slide, ax, y + 1.9, 0.68)

    x2 = MARGIN + cw + 0.95
    rect(slide, x2, y, cw, 3.85, fill=CARD_ALT, line=ORANGE, radius=0.12)
    icon_tile(slide, "agent", x2 + 0.35, y + 0.35, size=0.6, bg=ORANGE, fg=WHITE)
    text(slide, x2 + 0.35, y + 1.12, cw - 0.7, 0.3, "La solución: Nodi", size=19, bold=True)
    text(
        slide,
        x2 + 0.35,
        y + 1.55,
        cw - 0.7,
        0.6,
        "Un mentor que responde con fundamento y enseña el camino que recorrió.",
        size=13,
        color=MUTED,
        spacing=1.3,
    )
    feats = [("mic", "Texto y voz"), ("database", "RAG real"), ("stream", "Respuesta en vivo"), ("pipeline", "Flujo visible"), ("spark", "Infografías"), ("shield", "Solo temas tech")]
    fw = (cw - 0.7 - 0.2) / 2
    for i, (ic, label) in enumerate(feats):
        fx = x2 + 0.35 + (i % 2) * (fw + 0.2)
        fy = y + 2.35 + (i // 2) * 0.52
        place_icon(slide, ic, fx, fy + 0.04, 0.26, GREEN)
        text(slide, fx + 0.38, fy + 0.03, fw - 0.4, 0.28, label, size=12, bold=True)
    footer(slide)
    return slide


def slide_demo(prs):
    slide = new_slide(prs)
    header(slide, "02 · Demo end-to-end", "Cuatro momentos para demostrar el producto", "Funcionamiento real, datos reales e inteligencia con fundamento", number=4)

    cards = [
        ("1", "search", "RAG con acierto", "“¿Qué es Git?”", "Se activan la herramienta y el corpus; la respuesta cita fragmentos reales."),
        ("2", "database", "RAG sin resultado", "“¿Qué es ITIL?”", "La herramienta corre, no hay coincidencias y Nodi lo admite."),
        ("3", "shield", "Límite temático", "“Dame una receta”", "No llama la herramienta y redirige la conversación a tecnología."),
        ("4", "mic", "Interacción", "Micrófono o chip", "Tema relacionado como pregunta y avatar leyendo el resumen."),
    ]
    gap = 0.34
    cw = (CONTENT_W - 3 * gap) / 4
    y = 2.15
    for i, (n, ic, title, prompt, detail) in enumerate(cards):
        x = MARGIN + i * (cw + gap)
        rect(slide, x, y, cw, 3.1, fill=CARD, line=ORANGE if i == 0 else LINE, radius=0.12)
        icon_tile(slide, ic, x + 0.28, y + 0.3, size=0.56, bg=ORANGE if i == 0 else LINE, fg=WHITE if i == 0 else GREEN)
        text(slide, x + cw - 0.6, y + 0.36, 0.4, 0.3, n, size=17, color=LINE_SOFT, bold=True, align=PP_ALIGN.RIGHT)
        text(slide, x + 0.28, y + 1.05, cw - 0.56, 0.3, title, size=14.5, bold=True)
        rect(slide, x + 0.28, y + 1.45, cw - 0.56, 0.42, fill=NAVY_DEEP, radius=0.08)
        text(slide, x + 0.4, y + 1.55, cw - 0.8, 0.26, prompt, size=11.5, color=GREEN, bold=True)
        text(slide, x + 0.28, y + 2.02, cw - 0.56, 0.9, detail, size=11, color=MUTED, spacing=1.25)
        if i < 3:
            flow_arrow(slide, x + cw + 0.06, y + 1.55, gap - 0.12)

    y2 = 5.55
    rect(slide, MARGIN, y2, CONTENT_W, 1.0, fill=CARD_ALT, line=LINE, radius=0.1)
    icon_tile(slide, "play", MARGIN + 0.3, y2 + 0.22, size=0.56, bg=GREEN, fg=NAVY_DEEP)
    text(slide, MARGIN + 1.05, y2 + 0.22, 8.0, 0.28, "Momento principal: la aplicación publicada, en vivo", size=14.5, bold=True)
    text(slide, MARGIN + 1.05, y2 + 0.56, 8.0, 0.26, "tecno-chatbot-portal.vercel.app · no es una grabación", size=11.5, color=MUTED)
    pill(slide, W - MARGIN - 1.25, y2 + 0.33, "≈ 3 min", fill=ORANGE, height=0.34, size=12)
    footer(slide)
    return slide


def slide_architecture(prs):
    slide = new_slide(prs)
    header(slide, "03 · Arquitectura tecnológica", "De una pregunta a una explicación visual", "Tres fronteras claras: navegador, servidor Next.js y servicio de IA", number=5)

    def node(x, y, w, h, ic, title, detail, accent=False):
        rect(slide, x, y, w, h, fill=ORANGE if accent else CARD, line=ORANGE if accent else LINE, radius=0.1)
        place_icon(slide, ic, x + w / 2 - 0.16, y + 0.16, 0.32, WHITE if accent else GREEN)
        text(slide, x + 0.06, y + 0.56, w - 0.12, 0.24, title, size=11.5, bold=True, align=PP_ALIGN.CENTER)
        text(slide, x + 0.06, y + 0.82, w - 0.12, 0.22, detail, size=9, color=WHITE if accent else MUTED, align=PP_ALIGN.CENTER)

    nw, nh = 1.92, 1.16
    gap = (CONTENT_W - 5 * nw) / 4
    y = 2.05
    row1 = [
        ("user", "Usuario", "texto o voz", False),
        ("browser", "ChatApp", "React · Web Speech", False),
        ("code", "/api/chat", "Next.js · Node", False),
        ("model", "Claude Haiku", "razona y pide tool", True),
        ("tool", "1 herramienta", "consulta el corpus", False),
    ]
    for i, (ic, t, d, acc) in enumerate(row1):
        x = MARGIN + i * (nw + gap)
        node(x, y, nw, nh, ic, t, d, acc)
        if i < 4:
            flow_arrow(slide, x + nw + 0.05, y + nh / 2, gap - 0.1)

    xs = [MARGIN + i * (nw + gap) for i in range(5)]
    down_arrow(slide, xs[4] + nw / 2, y + nh + 0.06, 0.44)
    text(slide, xs[4] - 1.9, y + nh + 0.14, 1.75, 0.24, "la herramienta lee el corpus", size=9.5, color=MUTED, align=PP_ALIGN.RIGHT, italic=True)

    # Fila de retorno: se lee de derecha a izquierda, del corpus hasta la voz.
    y2 = y + nh + 0.56
    row2 = [
        ("agent", "Avatar y voz", "lee el resumen"),
        ("spark", "UI educativa", "Markdown · infografía"),
        ("stream", "NDJSON", "delta · done"),
        ("layers", "Fragmentos", "fuente y score"),
        ("database", "knowledge/*.md", "corpus local"),
    ]
    for i, (ic, t, d) in enumerate(row2):
        node(xs[i], y2, nw, nh, ic, t, d)
        if i < 4:
            flow_arrow(slide, xs[i] + nw + 0.05, y2 + nh / 2, gap - 0.1, color=GREEN, left=True)

    y3 = y2 + nh + 0.42
    layers = [
        ("browser", "Navegador", "Next.js · React · Tailwind · Web Speech API · avatar SVG"),
        ("code", "Servidor", "API Route · validación · tool use · RAG · streaming NDJSON"),
        ("model", "Anthropic", "Claude Haiku · system prompt · function calling"),
    ]
    lw = (CONTENT_W - 0.5) / 3
    for i, (ic, t, d) in enumerate(layers):
        x = MARGIN + i * (lw + 0.25)
        rect(slide, x, y3, lw, 0.92, fill=CARD_ALT, line=LINE, radius=0.1)
        icon_tile(slide, ic, x + 0.22, y3 + 0.22, size=0.48, bg=LINE, fg=GREEN)
        text(slide, x + 0.84, y3 + 0.18, lw - 1.0, 0.26, t, size=13, bold=True)
        text(slide, x + 0.84, y3 + 0.46, lw - 1.0, 0.4, d, size=10, color=MUTED, spacing=1.15)

    y4 = y3 + 1.06
    rect(slide, MARGIN, y4, CONTENT_W, 0.5, fill=NAVY_DEEP, line=ORANGE, radius=0.08)
    place_icon(slide, "key", MARGIN + 0.22, y4 + 0.12, 0.26, ORANGE)
    text(slide, MARGIN + 0.6, y4 + 0.13, CONTENT_W - 0.9, 0.26, "Frontera de seguridad: la API key vive solo en el servidor y en las variables protegidas de Vercel.", size=11.5, bold=True)
    return slide


def slide_intelligence(prs):
    slide = new_slide(prs)
    header(slide, "04 · Inteligencia con fundamento", "Claude propone; Node ejecuta", "Una sola herramienta, tal como pide el laboratorio", number=6)

    steps = [
        ("01", "model", "Prompt engineering", "Rol educativo, español claro, límite temático y formato de salida."),
        ("02", "tool", "Tool use", "Claude decide cuándo pedir consultar_conocimiento_tech."),
        ("03", "search", "RAG léxico", "Node busca los términos relevantes en el corpus Markdown local."),
        ("04", "spark", "Salida estructurada", "Respuesta, resumen hablado, temas relacionados e infografía."),
    ]
    lw = CONTENT_W * 0.53
    y = 2.1
    for i, (n, ic, title, detail) in enumerate(steps):
        cy = y + i * 1.08
        rect(slide, MARGIN, cy, lw, 0.86, fill=CARD, line=LINE, radius=0.1)
        text(slide, MARGIN + 0.24, cy + 0.25, 0.5, 0.34, n, size=17, color=ORANGE, bold=True)
        icon_tile(slide, ic, MARGIN + 0.8, cy + 0.18, size=0.5, bg=LINE, fg=GREEN)
        text(slide, MARGIN + 1.45, cy + 0.15, lw - 1.7, 0.26, title, size=13.5, bold=True)
        text(slide, MARGIN + 1.45, cy + 0.43, lw - 1.7, 0.4, detail, size=10.5, color=MUTED, spacing=1.15)
        if i < 3:
            down_arrow(slide, MARGIN + 0.49, cy + 0.88, 0.2, color=LINE_SOFT)

    rx = MARGIN + lw + 0.4
    rw = CONTENT_W - lw - 0.4
    rect(slide, rx, y, rw, 4.1, fill=CARD_ALT, line=ORANGE, radius=0.12)
    text(slide, rx + 0.3, y + 0.26, rw - 0.6, 0.26, "Contrato de la única herramienta", size=13.5, bold=True)
    rect(slide, rx + 0.3, y + 0.66, rw - 0.6, 0.86, fill=NAVY_DEEP, radius=0.08)
    text(slide, rx + 0.48, y + 0.8, rw - 0.96, 0.6, ["consultar_conocimiento_tech({", '  consulta: "¿Qué es Git?"', "})"], size=11.5, color=GREEN, spacing=1.15)
    down_arrow(slide, rx + rw / 2, y + 1.62, 0.36)
    rect(slide, rx + 0.3, y + 2.08, rw - 0.6, 0.78, fill=CARD, line=LINE, radius=0.08)
    text(slide, rx + 0.48, y + 2.22, rw - 0.96, 0.22, "Retorno", size=9.5, color=ORANGE, bold=True, caps=True, space=1.2)
    text(slide, rx + 0.48, y + 2.46, rw - 0.96, 0.26, "Hasta 4 fragmentos, con fuente y score", size=12, bold=True)
    rect(slide, rx + 0.3, y + 2.98, rw - 0.6, 1.05, fill=NAVY_DEEP, line=GREEN, radius=0.08)
    place_icon(slide, "shield", rx + 0.48, y + 3.14, 0.26, GREEN)
    text(slide, rx + 0.86, y + 3.14, rw - 1.2, 0.24, "Sin alucinar fuentes", size=12, bold=True, color=GREEN)
    text(slide, rx + 0.48, y + 3.46, rw - 0.96, 0.5, "Si no hay coincidencias, lo dice. Nunca finge una búsqueda en internet.", size=10.5, color=MUTED, spacing=1.2)
    footer(slide)
    return slide


def slide_quality(prs):
    slide = new_slide(prs)
    header(slide, "05 · Calidad verificada", "La confianza se construye por capas", "Cada nivel reduce un tipo distinto de riesgo antes de publicar", number=7)

    # Pirámide de pruebas.
    cx, top = 3.85, 2.2
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(cx - 0.85), Inches(top), Inches(1.7), Inches(0.95))
    tri.fill.solid()
    tri.fill.fore_color.rgb = rgb(ORANGE)
    tri.line.fill.background()
    tri.shadow.inherit = False
    text(slide, cx - 0.85, top + 0.62, 1.7, 0.24, "E2E", size=11.5, bold=True, align=PP_ALIGN.CENTER)

    mid = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(cx - 1.95), Inches(top + 1.05), Inches(3.9), Inches(1.0))
    mid.adjustments[0] = 0.28
    mid.fill.solid()
    mid.fill.fore_color.rgb = rgb(GREEN)
    mid.line.fill.background()
    mid.shadow.inherit = False
    text(slide, cx - 1.95, top + 1.48, 3.9, 0.26, "INTEGRACIÓN", size=12, bold=True, color=NAVY_DEEP, align=PP_ALIGN.CENTER)

    base = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(cx - 3.1), Inches(top + 2.15), Inches(6.2), Inches(1.15))
    base.adjustments[0] = 0.19
    base.fill.solid()
    base.fill.fore_color.rgb = rgb(CARD)
    base.line.color.rgb = rgb(LINE_SOFT)
    base.line.width = Pt(1)
    base.shadow.inherit = False
    text(slide, cx - 3.1, top + 2.6, 6.2, 0.26, "UNITARIAS", size=12.5, bold=True, align=PP_ALIGN.CENTER)
    text(slide, cx - 3.1, top + 2.88, 6.2, 0.24, "parseo · validación · búsqueda léxica", size=10, color=MUTED, align=PP_ALIGN.CENTER)

    text(slide, cx - 3.4, top + 3.55, 6.8, 0.24, "Base ancha, cima estrecha: muchas pruebas rápidas y pocas de navegador.", size=10.5, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    # Detalle a la derecha.
    rx = 7.55
    rw = W - MARGIN - rx
    cards = [
        ("code", "Pruebas de lógica", "Parseo de la respuesta e infografías, scoring léxico, RAG y validación del request."),
        ("browser", "Prueba E2E con Playwright", "Abre Chromium, pregunta “¿Qué es Git?”, verifica herramienta, RAG e infografía, y encadena un tema relacionado."),
    ]
    y = 2.2
    for ic, title, detail in cards:
        rect(slide, rx, y, rw, 1.28, fill=CARD, line=LINE, radius=0.1)
        icon_tile(slide, ic, rx + 0.26, y + 0.26, size=0.5, bg=LINE, fg=GREEN)
        text(slide, rx + 0.9, y + 0.24, rw - 1.15, 0.26, title, size=13, bold=True)
        text(slide, rx + 0.9, y + 0.54, rw - 1.15, 0.66, detail, size=10.5, color=MUTED, spacing=1.2)
        y += 1.45

    gw = (rw - 3 * 0.14) / 4
    for i, label in enumerate(["Tests", "Tipos", "Lint", "Build"]):
        gx = rx + i * (gw + 0.14)
        rect(slide, gx, y, gw, 0.78, fill=CARD_ALT, line=LINE, radius=0.09)
        place_icon(slide, "check", gx + gw / 2 - 0.13, y + 0.13, 0.26, GREEN)
        text(slide, gx, y + 0.46, gw, 0.22, label, size=10.5, bold=True, align=PP_ALIGN.CENTER)

    y += 0.95
    rect(slide, rx, y, rw, 0.6, fill=NAVY_DEEP, line=GREEN, radius=0.08)
    text(slide, rx + 0.24, y + 0.18, rw - 0.48, 0.26, "Un fallo bloquea el pipeline y deja evidencia reproducible.", size=11, bold=True, color=GREEN)
    footer(slide)
    return slide


def slide_devops(prs):
    slide = new_slide(prs)
    header(slide, "06 · DevOps · CI/CD", "Un push activa dos caminos automáticos", "GitHub conserva y verifica el código; Vercel publica la aplicación", number=8)

    stages = [
        ("code", "Código", "commit local", False),
        ("git", "GitHub", "push a main", False),
        ("pipeline", "GitHub Actions", "verificación", True),
        ("deploy", "Vercel", "build y deploy", True),
        ("browser", "Producción", "URL pública", False),
    ]
    nw, nh = 1.92, 1.2
    gap = (CONTENT_W - 5 * nw) / 4
    y = 2.02
    for i, (ic, t, d, acc) in enumerate(stages):
        x = MARGIN + i * (nw + gap)
        rect(slide, x, y, nw, nh, fill=ORANGE if acc else CARD, line=ORANGE if acc else LINE, radius=0.1)
        place_icon(slide, ic, x + nw / 2 - 0.17, y + 0.17, 0.34, WHITE if acc else GREEN)
        text(slide, x + 0.06, y + 0.6, nw - 0.12, 0.24, t, size=12, bold=True, align=PP_ALIGN.CENTER)
        text(slide, x + 0.06, y + 0.87, nw - 0.12, 0.22, d, size=9.5, color=WHITE if acc else MUTED, align=PP_ALIGN.CENTER)
        if i < 4:
            flow_arrow(slide, x + nw + 0.05, y + nh / 2, gap - 0.1)

    y2 = y + nh + 0.4
    cw = (CONTENT_W - 0.35) / 2
    # CI
    rect(slide, MARGIN, y2, cw, 2.5, fill=CARD, line=LINE, radius=0.12)
    icon_tile(slide, "pipeline", MARGIN + 0.28, y2 + 0.24, size=0.5, bg=LINE, fg=GREEN)
    text(slide, MARGIN + 0.92, y2 + 0.24, cw - 2.0, 0.28, "GitHub Actions", size=14.5, bold=True)
    text(slide, MARGIN + 0.92, y2 + 0.54, cw - 2.0, 0.22, "Integración continua en cada push", size=10.5, color=MUTED)
    pill(slide, MARGIN + cw - 0.75, y2 + 0.28, "CI", fill=GREEN, color=NAVY_DEEP, size=11)
    steps = ["npm ci", "Pruebas unitarias e integración", "TypeScript y ESLint", "Build de producción", "Playwright en Chromium"]
    for i, step in enumerate(steps):
        sy = y2 + 0.95 + i * 0.29
        place_icon(slide, "check", MARGIN + 0.32, sy + 0.02, 0.19, GREEN)
        text(slide, MARGIN + 0.62, sy, cw - 0.9, 0.24, step, size=11, color=WHITE if i else WHITE)

    # CD
    x2 = MARGIN + cw + 0.35
    rect(slide, x2, y2, cw, 2.5, fill=CARD, line=LINE, radius=0.12)
    icon_tile(slide, "cloud", x2 + 0.28, y2 + 0.24, size=0.5, bg=LINE, fg=GREEN)
    text(slide, x2 + 0.92, y2 + 0.24, cw - 2.0, 0.28, "Vercel", size=14.5, bold=True)
    text(slide, x2 + 0.92, y2 + 0.54, cw - 2.0, 0.22, "Entrega continua sin pasos manuales", size=10.5, color=MUTED)
    pill(slide, x2 + cw - 0.78, y2 + 0.28, "CD", fill=ORANGE, size=11)
    cd = [("git", "Detecta el commit en main"), ("code", "Construye la app Next.js"), ("browser", "Publica la nueva versión")]
    for i, (ic, label) in enumerate(cd):
        sy = y2 + 0.95 + i * 0.54
        icon_tile(slide, ic, x2 + 0.32, sy, size=0.36, bg=NAVY_DEEP, fg=GREEN, radius=0.07)
        text(slide, x2 + 0.82, sy + 0.05, cw - 1.1, 0.26, label, size=11.5, bold=True)
        if i < 2:
            down_arrow(slide, x2 + 0.5, sy + 0.38, 0.16, color=LINE_SOFT)

    y3 = y2 + 2.66
    rect(slide, MARGIN, y3, CONTENT_W, 0.52, fill=NAVY_DEEP, line=ORANGE, radius=0.08)
    place_icon(slide, "layers", MARGIN + 0.22, y3 + 0.13, 0.26, ORANGE)
    text(slide, MARGIN + 0.6, y3 + 0.14, CONTENT_W - 0.9, 0.26, "Precisión: hoy CI y Vercel arrancan en paralelo desde el mismo push; ambos quedan como evidencia del despliegue.", size=11, bold=True)
    return slide


def slide_decisions(prs):
    slide = new_slide(prs)
    header(slide, "07 · Decisiones", "Menos piezas, más claridad", "Cada decisión prioriza aprendizaje, trazabilidad y alcance controlado", number=9)

    decisions = [
        ("tool", "Una sola herramienta", "Cumple el laboratorio y hace observable el contrato entre modelo y código.", "Sin orquestaciones innecesarias"),
        ("search", "RAG léxico", "Simple, local, explicable y sin costo adicional de infraestructura.", "Sin base vectorial ni embeddings"),
        ("stream", "Haiku con NDJSON", "Menor latencia percibida: el texto aparece mientras se genera.", "Sin esperar la respuesta completa"),
    ]
    cw = (CONTENT_W - 0.6) / 3
    y = 2.1
    for i, (ic, title, why, avoided) in enumerate(decisions):
        x = MARGIN + i * (cw + 0.3)
        rect(slide, x, y, cw, 2.6, fill=CARD, line=LINE, radius=0.12)
        icon_tile(slide, ic, x + 0.3, y + 0.3, size=0.6, bg=ORANGE, fg=WHITE)
        text(slide, x + cw - 0.85, y + 0.32, 0.55, 0.4, f"0{i + 1}", size=24, color=LINE_SOFT, bold=True, align=PP_ALIGN.RIGHT)
        text(slide, x + 0.3, y + 1.08, cw - 0.6, 0.3, title, size=15, bold=True)
        text(slide, x + 0.3, y + 1.48, cw - 0.6, 0.62, why, size=11.5, color=MUTED, spacing=1.25)
        rect(slide, x + 0.3, y + 2.14, cw - 0.6, 0.02, fill=LINE, radius=0)
        text(slide, x + 0.3, y + 2.24, cw - 0.6, 0.24, avoided, size=10, color=LINE_SOFT, italic=True)

    y2 = y + 2.85
    text(slide, MARGIN, y2, CONTENT_W, 0.26, "Cobertura de la rúbrica", size=11, color=ORANGE, bold=True, caps=True, space=1.2)
    labels = ["End-to-end", "Datos reales", "IA con fundamento", "Pruebas", "Publicado y vivo", "Criterio propio"]
    gw = (CONTENT_W - 5 * 0.16) / 6
    for i, label in enumerate(labels):
        gx = MARGIN + i * (gw + 0.16)
        rect(slide, gx, y2 + 0.34, gw, 0.78, fill=CARD_ALT, line=GREEN, radius=0.09)
        place_icon(slide, "check", gx + gw / 2 - 0.13, y2 + 0.46, 0.26, GREEN)
        text(slide, gx + 0.05, y2 + 0.78, gw - 0.1, 0.24, label, size=10, bold=True, align=PP_ALIGN.CENTER)
    footer(slide)
    return slide


def slide_closing(prs):
    slide = new_slide(prs, NAVY_DEEP)
    rect(slide, 0, 0, 0.16, H, fill=ORANGE, radius=0)
    text(slide, 1.0, 1.15, 7, 0.24, "Cierre", size=11, color=ORANGE, bold=True, caps=True, space=1.6)
    text(slide, 0.96, 1.5, 6.6, 1.3, "Nodi convierte arquitectura en aprendizaje", size=38, bold=True, spacing=0.98)
    text(
        slide,
        1.0,
        3.0,
        6.3,
        0.8,
        "No solo responde preguntas: muestra el recorrido de la información, "
        "fundamenta lo que dice y demuestra que fue probado y publicado.",
        size=14,
        color=MUTED,
        spacing=1.3,
    )
    rect(slide, 1.0, 4.02, 0.05, 1.62, fill=ORANGE, radius=0)
    text(slide, 1.3, 4.02, 5.5, 1.62, ["Una función.", "Un corpus.", "Un stream.", "Un push."], size=21, bold=True, spacing=1.18)
    x = 1.0
    for label, fill in [("Aplicación publicada", GREEN), ("CI en verde", None), ("Repositorio en GitHub", None)]:
        x += pill(slide, x, 6.05, label, fill=fill, color=NAVY_DEEP if fill else WHITE, height=0.34, size=11) + 0.14
    text(slide, 1.0, 6.62, 7, 0.24, "tecno-chatbot-portal.vercel.app", size=11.5, color=MUTED)

    # Ciclo: pregunta, conocimiento, respuesta, repositorio, despliegue.
    cx, cy, r = 10.4, 3.75, 1.85
    labels = [("user", "Pregunta"), ("database", "Corpus"), ("stream", "Respuesta"), ("git", "Repositorio"), ("deploy", "Despliegue")]
    pts = []
    for i in range(5):
        a = math.radians(-90 + i * 72)
        pts.append((cx + r * math.cos(a), cy + r * math.sin(a)))
    for i, (px, py) in enumerate(pts):
        nx, ny = pts[(i + 1) % 5]
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(px), Inches(py), Inches(nx), Inches(ny))
        c.line.color.rgb = rgb(LINE)
        c.line.width = Pt(1.5)
    d = 1.5
    rect(slide, cx - d / 2, cy - d / 2, d, d, fill=ORANGE, radius=d / 2)
    place_icon(slide, "agent", cx - 0.38, cy - 0.38, 0.76, WHITE)
    for (px, py), (name, label) in zip(pts, labels):
        s = 0.8
        rect(slide, px - s / 2, py - s / 2, s, s, fill=CARD, line=LINE_SOFT, radius=s / 2)
        place_icon(slide, name, px - 0.19, py - 0.19, 0.38, GREEN)
        text(slide, px - 0.9, py + 0.46, 1.8, 0.22, label, size=9.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    slide_cover(prs)
    slide_agenda(prs)
    slide_challenge(prs)
    slide_demo(prs)
    slide_architecture(prs)
    slide_intelligence(prs)
    slide_quality(prs)
    slide_devops(prs)
    slide_decisions(prs)
    slide_closing(prs)

    out = Path(__file__).resolve().parent / "nodi-presentacion.pptx"
    prs.save(out)
    shutil.rmtree(_tmp, ignore_errors=True)
    print(f"{out}  ·  {len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas")


if __name__ == "__main__":
    main()
